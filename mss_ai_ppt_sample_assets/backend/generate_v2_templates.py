"""
Script to generate V2 PPTX templates for MSS reports.
All content is placeholders - NO hardcoded text.

Design style based on professional security report templates:
- Primary color: #0A4275 (deep professional blue)
- Clean card-based layout (rectangular, no shadows)
- Large bold numbers for KPIs
- Clear visual hierarchy with section headers

Run this script to create:
- mss_executive_v2.pptx (8 slides, management version)
- mss_technical_v2.pptx (10 slides, technical version)
"""

from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ============================================================
# Color Palette (based on professional security report style)
# ============================================================

# Primary brand color - deep professional blue
PRIMARY_BLUE = RGBColor(10, 66, 117)       # #0A4275 - Main brand color
PRIMARY_BLUE_LIGHT = RGBColor(49, 130, 206) # #3182CE - Accent blue

# Text colors
TEXT_DARK = RGBColor(51, 51, 51)           # #333333 - Primary text
TEXT_MEDIUM = RGBColor(85, 85, 85)         # #555555 - Secondary text
TEXT_LIGHT = RGBColor(102, 102, 102)       # #666666 - Muted text
TEXT_WHITE = RGBColor(255, 255, 255)       # #FFFFFF - White text

# Background colors
BG_WHITE = RGBColor(255, 255, 255)         # #FFFFFF - Pure white
BG_LIGHT_GRAY = RGBColor(238, 238, 238)    # #EEEEEE - Light gray panel
BG_SECTION = RGBColor(226, 232, 240)       # #E2E8F0 - Section background
BG_CARD = RGBColor(237, 242, 247)          # #EDF2F7 - Card background

# Status colors
STATUS_GREEN = RGBColor(0, 176, 80)        # #00B050 - Success/positive
STATUS_GREEN_ALT = RGBColor(56, 161, 105)  # #38A169 - Alternative green
STATUS_RED = RGBColor(220, 53, 69)         # #DC3545 - Alert/danger
STATUS_RED_ALT = RGBColor(220, 38, 38)     # #DC2626 - Critical
STATUS_ORANGE = RGBColor(234, 88, 12)      # #EA580C - Warning

# Legacy compatibility aliases
LIGHT_PRIMARY = PRIMARY_BLUE
LIGHT_ACCENT = PRIMARY_BLUE_LIGHT
LIGHT_BG = BG_WHITE
LIGHT_PANEL = BG_WHITE
LIGHT_BORDER = BG_SECTION
LIGHT_MUTED = TEXT_LIGHT

DARK_PRIMARY = RGBColor(15, 23, 42)        # #0F172A - Dark slate
DARK_ACCENT = RGBColor(34, 197, 94)        # #22C55E - Green
DARK_BG = RGBColor(11, 18, 32)             # #0B1220
DARK_PANEL = RGBColor(17, 24, 39)          # #111827
DARK_BORDER = RGBColor(31, 41, 55)         # #1F2937
DARK_TEXT = RGBColor(226, 232, 240)        # #E2E8F0
DARK_MUTED = RGBColor(148, 163, 184)       # #94A3B8


def add_placeholder_box(slide, left, top, width, height, token, font_size=12, bold=False, color=None):
    """Add a placeholder text box with {{TOKEN}} format."""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{{{{{token}}}}}"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    return shape


def _set_slide_bg(slide, color: RGBColor) -> None:
    """Set slide background color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: RGBColor,
    line_color: Optional[RGBColor] = None,
    transparency: Optional[float] = None,
):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if transparency is not None:
        shape.fill.transparency = max(0.0, min(1.0, transparency))

    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def _style_textbox(shape, font_color: RGBColor, font_size: Optional[int] = None, bold: Optional[bool] = None, align: Optional[PP_ALIGN] = None):
    """Style a text box with consistent formatting."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for p in tf.paragraphs:
        if font_size is not None:
            p.font.size = Pt(font_size)
        if bold is not None:
            p.font.bold = bold
        p.font.color.rgb = font_color
        p.font.name = "Microsoft YaHei"
        if align is not None:
            p.alignment = align


def _add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
    border: Optional[RGBColor] = None,
):
    """Add a card container (rectangular, no shadow)."""
    # Use rectangle instead of rounded rectangle, no shadow
    return _add_rect(slide, left, top, width, height, fill, line_color=border, transparency=0.0)


def _add_header_stripe(slide, slide_width: float = 13.333):
    """Add the top header stripe like the reference template."""
    # Top accent line - thin primary color stripe
    _add_rect(slide, 0, 0, slide_width, 0.08, PRIMARY_BLUE, transparency=0.0)


def _add_section_header(slide, left: float, top: float, width: float, token: str,
                        height: float = 0.5, font_size: int = 18, with_icon: bool = True):
    """Add a section header with optional icon decoration."""
    # Left accent bar (thinner)
    _add_rect(slide, left, top, 0.05, height, PRIMARY_BLUE, transparency=0.0)

    # Header text
    shape = add_placeholder_box(slide, left + 0.15, top, width - 0.15, height,
                                token, font_size=font_size, bold=True, color=PRIMARY_BLUE)
    _style_textbox(shape, PRIMARY_BLUE, font_size=font_size, bold=True)
    return shape


def _add_kpi_card(slide, left: float, top: float, width: float, height: float,
                  value_token: str, label_token: Optional[str] = None,
                  value_size: int = 36, label_size: int = 12):
    """Add a KPI display card with large number and label."""
    # Card background
    _add_card(slide, left, top, width, height, BG_CARD, border=BG_SECTION)

    # Large value number
    value_top = top + 0.15 if label_token else top + (height - 0.6) / 2
    value_shape = add_placeholder_box(slide, left + 0.1, value_top, width - 0.2, 0.6,
                                      value_token, font_size=value_size, bold=True, color=PRIMARY_BLUE)
    _style_textbox(value_shape, PRIMARY_BLUE, font_size=value_size, bold=True, align=PP_ALIGN.CENTER)

    # Label below
    if label_token:
        label_shape = add_placeholder_box(slide, left + 0.1, top + height - 0.4, width - 0.2, 0.35,
                                          label_token, font_size=label_size, color=TEXT_MEDIUM)
        _style_textbox(label_shape, TEXT_MEDIUM, font_size=label_size, align=PP_ALIGN.CENTER)

    return value_shape


def _light_base(slide):
    """Set up light theme base with header stripe."""
    _set_slide_bg(slide, BG_WHITE)
    _add_header_stripe(slide)


def _dark_base(slide):
    """Set up dark theme base."""
    _set_slide_bg(slide, DARK_BG)
    # Top accent line
    _add_rect(slide, 0, 0, 13.333, 0.08, DARK_ACCENT, transparency=0.0)


def create_executive_template():
    """Create the management/executive version template (8 slides).

    Design style based on professional security report:
    - Clean white background with blue accent stripe at top
    - Large bold numbers in primary blue for KPIs
    - Card-based layout with subtle borders
    - Clear section headers with left accent bars
    - Professional typography hierarchy
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ================================================================
    # Slide 1: Cover Page
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_WHITE)

    # Top accent stripe
    _add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY_BLUE, transparency=0.0)

    # Main title area
    title = add_placeholder_box(slide, 0.8, 2.8, 11.7, 1.0, "REPORT_TITLE",
                                font_size=42, bold=True, color=PRIMARY_BLUE)
    _style_textbox(title, PRIMARY_BLUE, font_size=42, bold=True, align=PP_ALIGN.CENTER)

    # Decorative line under title
    _add_rect(slide, 5.5, 3.9, 2.3, 0.05, PRIMARY_BLUE, transparency=0.0)

    # Customer and period info
    customer = add_placeholder_box(slide, 0.8, 4.3, 11.7, 0.5, "CUSTOMER_LABEL",
                                   font_size=18, color=TEXT_DARK)
    _style_textbox(customer, TEXT_DARK, font_size=18, align=PP_ALIGN.CENTER)

    period = add_placeholder_box(slide, 0.8, 4.85, 11.7, 0.45, "PERIOD_LABEL",
                                 font_size=16, color=TEXT_MEDIUM)
    _style_textbox(period, TEXT_MEDIUM, font_size=16, align=PP_ALIGN.CENTER)

    # Footer info
    conf = add_placeholder_box(slide, 0.8, 6.4, 5.5, 0.35, "CONFIDENTIALITY_LABEL",
                               font_size=14, color=TEXT_LIGHT)
    _style_textbox(conf, TEXT_LIGHT, font_size=14)

    gen = add_placeholder_box(slide, 6.5, 6.4, 6.0, 0.35, "GENERATED_AT_LABEL",
                              font_size=14, color=TEXT_LIGHT)
    _style_textbox(gen, TEXT_LIGHT, font_size=14, align=PP_ALIGN.RIGHT)

    # ================================================================
    # Slide 2: Executive Summary / Security Overview
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _light_base(slide)

    # Page title
    _add_section_header(slide, 0.6, 0.25, 12.1, "SLIDE_TITLE", height=0.55, font_size=27)

    # Headline card - prominent summary
    _add_card(slide, 0.6, 0.95, 12.1, 0.85, BG_CARD, border=BG_SECTION)
    hl = add_placeholder_box(slide, 0.8, 1.05, 11.7, 0.65, "HEADLINE",
                             font_size=18, bold=True, color=PRIMARY_BLUE)
    _style_textbox(hl, PRIMARY_BLUE, font_size=18, bold=True)

    # KPI Section - large numbers display
    _add_card(slide, 0.6, 1.95, 12.1, 1.3, BG_WHITE, border=BG_SECTION)
    kpi = add_placeholder_box(slide, 0.8, 2.1, 11.7, 1.0, "KPI_SECTION",
                              font_size=14, color=TEXT_DARK)
    _style_textbox(kpi, TEXT_DARK, font_size=14)

    # Summary paragraph
    _add_card(slide, 0.6, 3.4, 12.1, 1.6, BG_WHITE, border=BG_SECTION)
    sp = add_placeholder_box(slide, 0.8, 3.55, 11.7, 1.35, "SUMMARY_PARAGRAPH",
                             font_size=14, color=TEXT_DARK)
    _style_textbox(sp, TEXT_DARK, font_size=14)

    # Key Insights section
    _add_card(slide, 0.6, 5.15, 12.1, 2.15, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 5.15, 0.05, 2.15, PRIMARY_BLUE)  # Accent bar - full height
    kit = add_placeholder_box(slide, 0.8, 5.2, 11.5, 0.45, "KEY_INSIGHTS_TITLE",
                              font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(kit, PRIMARY_BLUE, font_size=16, bold=True)
    ki = add_placeholder_box(slide, 0.8, 5.7, 11.5, 1.5, "KEY_INSIGHTS",
                             font_size=13, color=TEXT_DARK)
    _style_textbox(ki, TEXT_DARK, font_size=13)

    # ================================================================
    # Slide 3: Alert Trends Analysis
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _light_base(slide)

    _add_section_header(slide, 0.6, 0.25, 12.1, "SLIDE_TITLE", height=0.55, font_size=27)

    # Left card - Trend section
    _add_card(slide, 0.6, 0.95, 5.85, 2.4, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 0.95, 0.05, 2.4, PRIMARY_BLUE)
    tst = add_placeholder_box(slide, 0.8, 1.0, 5.45, 0.4, "TREND_SECTION_TITLE",
                              font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(tst, PRIMARY_BLUE, font_size=16, bold=True)
    ta = add_placeholder_box(slide, 0.8, 1.45, 5.45, 1.8, "TREND_ANALYSIS",
                             font_size=12, color=TEXT_DARK)
    _style_textbox(ta, TEXT_DARK, font_size=12)

    # Right card - Categories
    _add_card(slide, 6.65, 0.95, 6.05, 2.4, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 6.65, 0.95, 0.05, 2.4, PRIMARY_BLUE)
    tct = add_placeholder_box(slide, 6.85, 1.0, 5.65, 0.4, "TOP_CATEGORIES_TITLE",
                              font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(tct, PRIMARY_BLUE, font_size=16, bold=True)
    tcl = add_placeholder_box(slide, 6.85, 1.45, 5.65, 1.8, "TOP_CATEGORIES_LIST",
                              font_size=12, color=TEXT_DARK)
    _style_textbox(tcl, TEXT_DARK, font_size=12)

    # Categories Insight
    _add_card(slide, 0.6, 3.5, 12.1, 1.65, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 3.5, 0.05, 1.65, PRIMARY_BLUE)
    cit = add_placeholder_box(slide, 0.8, 3.55, 11.7, 0.4, "CATEGORIES_INSIGHT_TITLE",
                              font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(cit, PRIMARY_BLUE, font_size=16, bold=True)
    tci = add_placeholder_box(slide, 0.8, 4.0, 11.7, 1.05, "TOP_CATEGORIES_INSIGHT",
                              font_size=13, color=TEXT_DARK)
    _style_textbox(tci, TEXT_DARK, font_size=13)

    # Month-over-Month comparison
    _add_card(slide, 0.6, 5.3, 12.1, 1.95, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 5.3, 0.05, 1.95, PRIMARY_BLUE)
    mt = add_placeholder_box(slide, 0.8, 5.35, 11.7, 0.4, "MOM_TITLE",
                             font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(mt, PRIMARY_BLUE, font_size=16, bold=True)
    mc = add_placeholder_box(slide, 0.8, 5.8, 11.7, 1.35, "MOM_COMPARISON",
                             font_size=13, color=TEXT_DARK)
    _style_textbox(mc, TEXT_DARK, font_size=13)

    # ================================================================
    # Slide 4: Major Incidents
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _light_base(slide)

    _add_section_header(slide, 0.6, 0.25, 12.1, "SLIDE_TITLE", height=0.55, font_size=27)

    # Incident Overview
    _add_card(slide, 0.6, 0.95, 12.1, 1.35, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 0.95, 0.05, 1.35, PRIMARY_BLUE)
    iot = add_placeholder_box(slide, 0.8, 1.0, 11.7, 0.4, "INCIDENT_OVERVIEW_TITLE",
                              font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(iot, PRIMARY_BLUE, font_size=16, bold=True)
    isum = add_placeholder_box(slide, 0.8, 1.45, 11.7, 0.75, "INCIDENT_SUMMARY",
                               font_size=13, color=TEXT_DARK)
    _style_textbox(isum, TEXT_DARK, font_size=13)

    # Incident Details / Table area
    _add_card(slide, 0.6, 2.45, 12.1, 3.15, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 2.45, 0.05, 3.15, PRIMARY_BLUE)
    idt = add_placeholder_box(slide, 0.8, 2.5, 11.7, 0.4, "INCIDENT_DETAILS_TITLE",
                              font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(idt, PRIMARY_BLUE, font_size=16, bold=True)
    idet = add_placeholder_box(slide, 0.8, 2.95, 11.7, 2.55, "INCIDENT_DETAILS",
                               font_size=11, color=TEXT_DARK)
    _style_textbox(idet, TEXT_DARK, font_size=11)

    # Incident Insight
    _add_card(slide, 0.6, 5.75, 12.1, 1.5, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 5.75, 0.05, 1.5, PRIMARY_BLUE)
    iit = add_placeholder_box(slide, 0.8, 5.8, 11.7, 0.4, "INCIDENT_INSIGHT_TITLE",
                              font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(iit, PRIMARY_BLUE, font_size=16, bold=True)
    ii = add_placeholder_box(slide, 0.8, 6.25, 11.7, 0.9, "INCIDENT_INSIGHT",
                             font_size=13, color=TEXT_DARK)
    _style_textbox(ii, TEXT_DARK, font_size=13)

    # ================================================================
    # Slide 5: Vulnerability & Exposure
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _light_base(slide)

    _add_section_header(slide, 0.6, 0.25, 12.1, "SLIDE_TITLE", height=0.55, font_size=27)

    # Vulnerability Stats - top banner
    _add_card(slide, 0.6, 0.95, 12.1, 0.95, BG_CARD, border=BG_SECTION)
    _add_rect(slide, 0.6, 0.95, 0.05, 0.95, PRIMARY_BLUE)
    vst = add_placeholder_box(slide, 0.8, 1.0, 2.5, 0.4, "VULN_STATS_TITLE",
                              font_size=14, bold=True, color=PRIMARY_BLUE)
    _style_textbox(vst, PRIMARY_BLUE, font_size=14, bold=True)
    vs = add_placeholder_box(slide, 3.3, 1.0, 9.2, 0.8, "VULN_STATS",
                             font_size=14, bold=True, color=TEXT_DARK)
    _style_textbox(vs, TEXT_DARK, font_size=14, bold=True)

    # Left - Vulnerability Overview
    _add_card(slide, 0.6, 2.05, 5.85, 1.55, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 2.05, 0.05, 1.55, PRIMARY_BLUE)
    vot = add_placeholder_box(slide, 0.8, 2.1, 5.45, 0.35, "VULN_OVERVIEW_TITLE",
                              font_size=14, bold=True, color=PRIMARY_BLUE)
    _style_textbox(vot, PRIMARY_BLUE, font_size=14, bold=True)
    vo = add_placeholder_box(slide, 0.8, 2.5, 5.45, 1.0, "VULN_OVERVIEW",
                             font_size=12, color=TEXT_DARK)
    _style_textbox(vo, TEXT_DARK, font_size=12)

    # Right - Exposure Stats
    _add_card(slide, 6.65, 2.05, 6.05, 1.55, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 6.65, 2.05, 0.05, 1.55, PRIMARY_BLUE)
    est = add_placeholder_box(slide, 6.85, 2.1, 5.65, 0.35, "EXPOSURE_STATS_TITLE",
                              font_size=14, bold=True, color=PRIMARY_BLUE)
    _style_textbox(est, PRIMARY_BLUE, font_size=14, bold=True)
    es = add_placeholder_box(slide, 6.85, 2.5, 5.65, 1.0, "EXPOSURE_STATS",
                             font_size=12, color=TEXT_DARK)
    _style_textbox(es, TEXT_DARK, font_size=12)

    # Top CVE List / Table
    _add_card(slide, 0.6, 3.75, 12.1, 1.85, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 3.75, 0.05, 1.85, PRIMARY_BLUE)
    tct = add_placeholder_box(slide, 0.8, 3.8, 11.7, 0.35, "TOP_CVE_TITLE",
                              font_size=14, bold=True, color=PRIMARY_BLUE)
    _style_textbox(tct, PRIMARY_BLUE, font_size=14, bold=True)
    tcl = add_placeholder_box(slide, 0.8, 4.2, 11.7, 1.3, "TOP_CVE_LIST",
                              font_size=11, color=TEXT_DARK)
    _style_textbox(tcl, TEXT_DARK, font_size=11)

    # CVE Analysis
    _add_card(slide, 0.6, 5.75, 12.1, 1.25, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 5.75, 0.05, 1.25, PRIMARY_BLUE)
    cat = add_placeholder_box(slide, 0.8, 5.8, 11.7, 0.35, "CVE_ANALYSIS_TITLE",
                              font_size=14, bold=True, color=PRIMARY_BLUE)
    _style_textbox(cat, PRIMARY_BLUE, font_size=14, bold=True)
    ana = add_placeholder_box(slide, 0.8, 6.2, 11.7, 0.7, "TOP_CVE_ANALYSIS",
                              font_size=12, color=TEXT_DARK)
    _style_textbox(ana, TEXT_DARK, font_size=12)

    # Exposure Summary footer
    exs = add_placeholder_box(slide, 0.6, 7.1, 12.1, 0.35, "EXPOSURE_SUMMARY",
                              font_size=10, color=TEXT_LIGHT)
    _style_textbox(exs, TEXT_LIGHT, font_size=10)

    # ================================================================
    # Slide 6: Cloud Security
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _light_base(slide)

    _add_section_header(slide, 0.6, 0.25, 12.1, "SLIDE_TITLE", height=0.55, font_size=27)

    # Cloud Accounts KPI card
    _add_card(slide, 0.6, 0.95, 3.2, 1.5, BG_CARD, border=BG_SECTION)
    cat = add_placeholder_box(slide, 0.75, 1.05, 2.9, 0.35, "CLOUD_ACCOUNTS_TITLE",
                              font_size=12, color=TEXT_MEDIUM)
    _style_textbox(cat, TEXT_MEDIUM, font_size=12, align=PP_ALIGN.CENTER)
    cac = add_placeholder_box(slide, 0.75, 1.45, 2.9, 0.9, "CLOUD_ACCOUNTS_COUNT",
                              font_size=36, bold=True, color=PRIMARY_BLUE)
    _style_textbox(cac, PRIMARY_BLUE, font_size=36, bold=True, align=PP_ALIGN.CENTER)

    # Cloud Risk section
    _add_card(slide, 4.0, 0.95, 8.7, 3.65, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 4.0, 0.95, 0.05, 3.65, PRIMARY_BLUE)
    crt = add_placeholder_box(slide, 4.2, 1.0, 8.3, 0.4, "CLOUD_RISK_TITLE",
                              font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(crt, PRIMARY_BLUE, font_size=16, bold=True)

    # Risk list (left half)
    crl = add_placeholder_box(slide, 4.2, 1.5, 4.0, 3.0, "CLOUD_RISK_LIST",
                              font_size=12, color=TEXT_DARK)
    _style_textbox(crl, TEXT_DARK, font_size=12)

    # Risk analysis (right half)
    _add_rect(slide, 8.3, 1.5, 0.02, 2.8, BG_SECTION)  # Divider line
    cat2 = add_placeholder_box(slide, 8.5, 1.5, 4.0, 0.35, "CLOUD_ANALYSIS_TITLE",
                               font_size=14, bold=True, color=PRIMARY_BLUE)
    _style_textbox(cat2, PRIMARY_BLUE, font_size=14, bold=True)
    crs = add_placeholder_box(slide, 8.5, 1.9, 4.0, 2.6, "CLOUD_RISK_SUMMARY",
                              font_size=12, color=TEXT_DARK)
    _style_textbox(crs, TEXT_DARK, font_size=12)

    # Cloud Recommendations
    _add_card(slide, 0.6, 4.75, 12.1, 2.5, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 4.75, 0.05, 2.5, PRIMARY_BLUE)
    crt2 = add_placeholder_box(slide, 0.8, 4.8, 11.7, 0.4, "CLOUD_REC_TITLE",
                               font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(crt2, PRIMARY_BLUE, font_size=16, bold=True)
    rec = add_placeholder_box(slide, 0.8, 5.25, 11.7, 1.9, "CLOUD_RECOMMENDATIONS",
                              font_size=13, color=TEXT_DARK)
    _style_textbox(rec, TEXT_DARK, font_size=13)

    # ================================================================
    # Slide 7: Recommendations & Action Plan
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _light_base(slide)

    _add_section_header(slide, 0.6, 0.25, 12.1, "SLIDE_TITLE", height=0.55, font_size=27)

    # P0 Actions (left)
    _add_card(slide, 0.6, 0.95, 5.85, 3.25, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 0.95, 0.05, 3.25, STATUS_RED)  # Red accent for urgent
    p0t = add_placeholder_box(slide, 0.8, 1.0, 5.45, 0.4, "P0_TITLE",
                              font_size=16, bold=True, color=STATUS_RED)
    _style_textbox(p0t, STATUS_RED, font_size=16, bold=True)
    p0a = add_placeholder_box(slide, 0.8, 1.45, 5.45, 2.65, "P0_ACTIONS",
                              font_size=12, color=TEXT_DARK)
    _style_textbox(p0a, TEXT_DARK, font_size=12)

    # P1 Actions (right)
    _add_card(slide, 6.65, 0.95, 6.05, 3.25, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 6.65, 0.95, 0.05, 3.25, STATUS_ORANGE)  # Orange for important
    p1t = add_placeholder_box(slide, 6.85, 1.0, 5.65, 0.4, "P1_TITLE",
                              font_size=16, bold=True, color=STATUS_ORANGE)
    _style_textbox(p1t, STATUS_ORANGE, font_size=16, bold=True)
    p1a = add_placeholder_box(slide, 6.85, 1.45, 5.65, 2.65, "P1_ACTIONS",
                              font_size=12, color=TEXT_DARK)
    _style_textbox(p1a, TEXT_DARK, font_size=12)

    # Strategic Recommendations
    _add_card(slide, 0.6, 4.35, 12.1, 2.9, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 4.35, 0.05, 2.9, PRIMARY_BLUE)
    srt = add_placeholder_box(slide, 0.8, 4.4, 11.7, 0.4, "STRATEGIC_TITLE",
                              font_size=16, bold=True, color=PRIMARY_BLUE)
    _style_textbox(srt, PRIMARY_BLUE, font_size=16, bold=True)
    srr = add_placeholder_box(slide, 0.8, 4.85, 11.7, 2.3, "STRATEGIC_RECOMMENDATIONS",
                              font_size=13, color=TEXT_DARK)
    _style_textbox(srr, TEXT_DARK, font_size=13)

    # ================================================================
    # Slide 8: Appendix & Notes
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _light_base(slide)

    _add_section_header(slide, 0.6, 0.25, 12.1, "SLIDE_TITLE", height=0.55, font_size=27)

    # Data Scope (left)
    _add_card(slide, 0.6, 0.95, 5.85, 1.85, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 0.95, 0.05, 1.85, PRIMARY_BLUE)
    dst = add_placeholder_box(slide, 0.8, 1.0, 5.45, 0.35, "DATA_SCOPE_TITLE",
                              font_size=14, bold=True, color=PRIMARY_BLUE)
    _style_textbox(dst, PRIMARY_BLUE, font_size=14, bold=True)
    ds = add_placeholder_box(slide, 0.8, 1.4, 5.45, 1.3, "DATA_SCOPE",
                             font_size=11, color=TEXT_DARK)
    _style_textbox(ds, TEXT_DARK, font_size=11)

    # Asset Coverage (right)
    _add_card(slide, 6.65, 0.95, 6.05, 1.85, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 6.65, 0.95, 0.05, 1.85, PRIMARY_BLUE)
    act = add_placeholder_box(slide, 6.85, 1.0, 5.65, 0.35, "ASSET_COVERAGE_TITLE",
                              font_size=14, bold=True, color=PRIMARY_BLUE)
    _style_textbox(act, PRIMARY_BLUE, font_size=14, bold=True)
    ac = add_placeholder_box(slide, 6.85, 1.4, 5.65, 1.3, "ASSET_COVERAGE",
                             font_size=11, color=TEXT_DARK)
    _style_textbox(ac, TEXT_DARK, font_size=11)

    # SLA Notes
    _add_card(slide, 0.6, 2.95, 12.1, 0.9, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 2.95, 0.05, 0.9, PRIMARY_BLUE)
    slat = add_placeholder_box(slide, 0.8, 3.0, 2.5, 0.35, "SLA_TITLE",
                               font_size=14, bold=True, color=PRIMARY_BLUE)
    _style_textbox(slat, PRIMARY_BLUE, font_size=14, bold=True)
    slan = add_placeholder_box(slide, 3.4, 3.0, 9.1, 0.75, "SLA_NOTES",
                               font_size=11, color=TEXT_DARK)
    _style_textbox(slan, TEXT_DARK, font_size=11)

    # Terminology
    _add_card(slide, 0.6, 4.0, 12.1, 2.65, BG_WHITE, border=BG_SECTION)
    _add_rect(slide, 0.6, 4.0, 0.05, 2.65, PRIMARY_BLUE)
    tt = add_placeholder_box(slide, 0.8, 4.05, 11.7, 0.35, "TERMINOLOGY_TITLE",
                             font_size=14, bold=True, color=PRIMARY_BLUE)
    _style_textbox(tt, PRIMARY_BLUE, font_size=14, bold=True)
    term = add_placeholder_box(slide, 0.8, 4.45, 11.7, 2.1, "TERMINOLOGY",
                               font_size=11, color=TEXT_DARK)
    _style_textbox(term, TEXT_DARK, font_size=11)

    # Contact Info footer
    ci = add_placeholder_box(slide, 0.6, 6.8, 12.1, 0.4, "CONTACT_INFO",
                             font_size=11, color=TEXT_LIGHT)
    _style_textbox(ci, TEXT_LIGHT, font_size=11, align=PP_ALIGN.CENTER)

    return prs


def create_technical_template():
    """Create the technical version template (10 slides).
    ALL content is placeholders - no hardcoded text.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_base(slide)
    _add_rect(slide, 0, 0, 0.18, 7.5, DARK_ACCENT, transparency=0.0)
    _add_card(slide, 0.65, 2.0, 12.3, 2.35, DARK_PANEL, DARK_BORDER)
    rt = add_placeholder_box(slide, 0.85, 2.2, 11.9, 0.95, "REPORT_TITLE", font_size=34, bold=True, color=DARK_ACCENT)
    _style_textbox(rt, DARK_ACCENT, font_size=34, bold=True)

    _add_card(slide, 0.65, 4.55, 7.8, 2.25, DARK_PANEL, DARK_BORDER)
    cl = add_placeholder_box(slide, 0.85, 4.72, 7.4, 0.55, "CUSTOMER_LABEL", font_size=18, bold=True, color=DARK_TEXT)
    _style_textbox(cl, DARK_TEXT, font_size=18, bold=True)
    pl = add_placeholder_box(slide, 0.85, 5.33, 7.4, 0.45, "PERIOD_LABEL", font_size=14, color=DARK_MUTED)
    _style_textbox(pl, DARK_MUTED, font_size=14)
    vt = add_placeholder_box(slide, 0.85, 5.83, 7.4, 0.45, "VERSION_TAG", font_size=12, bold=True, color=DARK_ACCENT)
    _style_textbox(vt, DARK_ACCENT, font_size=12, bold=True)
    cf = add_placeholder_box(slide, 0.85, 6.28, 7.4, 0.45, "CONFIDENTIALITY_LABEL", font_size=12, color=DARK_MUTED)
    _style_textbox(cf, DARK_MUTED, font_size=12)
    ga = add_placeholder_box(slide, 0.85, 6.73, 7.4, 0.45, "GENERATED_AT_LABEL", font_size=10, color=DARK_MUTED)
    _style_textbox(ga, DARK_MUTED, font_size=10)

    # Slide 2: Security Dashboard
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_base(slide)
    st = add_placeholder_box(slide, 0.65, 0.22, 12.3, 0.68, "SLIDE_TITLE", font_size=28, bold=True, color=DARK_ACCENT)
    _style_textbox(st, DARK_TEXT, font_size=28, bold=True)

    _add_card(slide, 0.65, 1.05, 12.3, 3.65, DARK_PANEL, DARK_BORDER)
    kpi = add_placeholder_box(slide, 0.85, 1.18, 11.9, 3.4, "KPI_DASHBOARD", font_size=11, color=DARK_TEXT)
    _style_textbox(kpi, DARK_TEXT, font_size=11)

    _add_card(slide, 0.65, 4.85, 12.3, 2.35, DARK_PANEL, DARK_BORDER)
    at = add_placeholder_box(slide, 0.85, 4.95, 11.9, 0.4, "ASSESSMENT_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(at, DARK_ACCENT, font_size=14, bold=True)
    ta = add_placeholder_box(slide, 0.85, 5.35, 11.9, 1.75, "TECHNICAL_ASSESSMENT", font_size=11, color=DARK_TEXT)
    _style_textbox(ta, DARK_TEXT, font_size=11)

    # Slide 3: Alert Deep Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_base(slide)
    st = add_placeholder_box(slide, 0.65, 0.22, 12.3, 0.68, "SLIDE_TITLE", font_size=28, bold=True, color=DARK_ACCENT)
    _style_textbox(st, DARK_TEXT, font_size=28, bold=True)

    _add_card(slide, 0.65, 1.05, 12.3, 1.35, DARK_PANEL, DARK_BORDER)
    svt = add_placeholder_box(slide, 0.85, 1.15, 11.9, 0.4, "SEVERITY_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(svt, DARK_ACCENT, font_size=14, bold=True)
    sv = add_placeholder_box(slide, 0.85, 1.55, 11.9, 0.75, "SEVERITY_BREAKDOWN", font_size=11, color=DARK_TEXT)
    _style_textbox(sv, DARK_TEXT, font_size=11)

    _add_card(slide, 0.65, 2.55, 6.05, 2.2, DARK_PANEL, DARK_BORDER)
    trt = add_placeholder_box(slide, 0.85, 2.65, 5.65, 0.4, "TOP_RULES_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(trt, DARK_ACCENT, font_size=14, bold=True)
    trtbl = add_placeholder_box(slide, 0.85, 3.05, 5.65, 1.6, "TOP_RULES_TABLE", font_size=9, color=DARK_TEXT)
    _style_textbox(trtbl, DARK_TEXT, font_size=9)

    _add_card(slide, 6.9, 2.55, 6.05, 2.2, DARK_PANEL, DARK_BORDER)
    rat = add_placeholder_box(slide, 7.1, 2.65, 5.65, 0.4, "RULES_ANALYSIS_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(rat, DARK_ACCENT, font_size=14, bold=True)
    raa = add_placeholder_box(slide, 7.1, 3.05, 5.65, 1.6, "TOP_RULES_ANALYSIS", font_size=10, color=DARK_TEXT)
    _style_textbox(raa, DARK_TEXT, font_size=10)

    _add_card(slide, 0.65, 4.9, 12.3, 2.3, DARK_PANEL, DARK_BORDER)
    fpt = add_placeholder_box(slide, 0.85, 5.0, 11.9, 0.4, "FP_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(fpt, DARK_ACCENT, font_size=14, bold=True)
    fpi = add_placeholder_box(slide, 0.85, 5.4, 11.9, 1.7, "FALSE_POSITIVE_INSIGHT", font_size=11, color=DARK_TEXT)
    _style_textbox(fpi, DARK_TEXT, font_size=11)

    # Slide 4: Incident Timeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_base(slide)
    st = add_placeholder_box(slide, 0.65, 0.22, 12.3, 0.68, "SLIDE_TITLE", font_size=28, bold=True, color=DARK_ACCENT)
    _style_textbox(st, DARK_TEXT, font_size=28, bold=True)

    _add_card(slide, 0.65, 1.05, 12.3, 2.95, DARK_PANEL, DARK_BORDER)
    tt = add_placeholder_box(slide, 0.85, 1.15, 11.9, 0.4, "TIMELINE_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(tt, DARK_ACCENT, font_size=14, bold=True)
    tn = add_placeholder_box(slide, 0.85, 1.55, 11.9, 2.35, "TIMELINE_NARRATIVE", font_size=11, color=DARK_TEXT)
    _style_textbox(tn, DARK_TEXT, font_size=11)

    _add_card(slide, 0.65, 4.15, 6.05, 3.05, DARK_PANEL, DARK_BORDER)
    rt = add_placeholder_box(slide, 0.85, 4.25, 5.65, 0.4, "RESPONSE_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(rt, DARK_ACCENT, font_size=14, bold=True)
    rm = add_placeholder_box(slide, 0.85, 4.65, 5.65, 2.45, "RESPONSE_METRICS", font_size=11, color=DARK_TEXT)
    _style_textbox(rm, DARK_TEXT, font_size=11)

    _add_card(slide, 6.9, 4.15, 6.05, 3.05, DARK_PANEL, DARK_BORDER)
    lt = add_placeholder_box(slide, 7.1, 4.25, 5.65, 0.4, "LESSONS_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(lt, DARK_ACCENT, font_size=14, bold=True)
    ll = add_placeholder_box(slide, 7.1, 4.65, 5.65, 2.45, "LESSONS_LEARNED", font_size=11, color=DARK_TEXT)
    _style_textbox(ll, DARK_TEXT, font_size=11)

    # Slide 5: Incident Details
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_base(slide)
    st = add_placeholder_box(slide, 0.65, 0.22, 12.3, 0.68, "SLIDE_TITLE", font_size=28, bold=True, color=DARK_ACCENT)
    _style_textbox(st, DARK_TEXT, font_size=28, bold=True)

    _add_card(slide, 0.65, 1.05, 12.3, 2.65, DARK_PANEL, DARK_BORDER)
    i1t = add_placeholder_box(slide, 0.85, 1.15, 11.9, 0.4, "INCIDENT1_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(i1t, DARK_ACCENT, font_size=14, bold=True)
    i1d = add_placeholder_box(slide, 0.85, 1.55, 11.9, 2.05, "INCIDENT_DETAIL_1", font_size=10, color=DARK_TEXT)
    _style_textbox(i1d, DARK_TEXT, font_size=10)

    _add_card(slide, 0.65, 3.85, 12.3, 2.65, DARK_PANEL, DARK_BORDER)
    i2t = add_placeholder_box(slide, 0.85, 3.95, 11.9, 0.4, "INCIDENT2_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(i2t, DARK_ACCENT, font_size=14, bold=True)
    i2d = add_placeholder_box(slide, 0.85, 4.35, 11.9, 2.05, "INCIDENT_DETAIL_2", font_size=10, color=DARK_TEXT)
    _style_textbox(i2d, DARK_TEXT, font_size=10)

    aps = add_placeholder_box(slide, 0.65, 6.95, 12.3, 0.45, "ATTACK_PATTERN_SUMMARY", font_size=9, color=DARK_MUTED)
    _style_textbox(aps, DARK_MUTED, font_size=9)

    # Slide 6: Vulnerability Details
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_base(slide)
    st = add_placeholder_box(slide, 0.65, 0.22, 12.3, 0.68, "SLIDE_TITLE", font_size=28, bold=True, color=DARK_ACCENT)
    _style_textbox(st, DARK_TEXT, font_size=28, bold=True)

    _add_card(slide, 0.65, 1.05, 12.3, 2.15, DARK_PANEL, DARK_BORDER)
    vdt = add_placeholder_box(slide, 0.85, 1.15, 11.9, 0.4, "VULN_DIST_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(vdt, DARK_ACCENT, font_size=14, bold=True)
    vd = add_placeholder_box(slide, 0.85, 1.55, 11.9, 1.55, "VULN_DISTRIBUTION", font_size=11, color=DARK_TEXT)
    _style_textbox(vd, DARK_TEXT, font_size=11)

    _add_card(slide, 0.65, 3.35, 12.3, 2.35, DARK_PANEL, DARK_BORDER)
    cat = add_placeholder_box(slide, 0.85, 3.45, 11.9, 0.4, "CVE_ANALYSIS_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(cat, DARK_ACCENT, font_size=14, bold=True)
    cta = add_placeholder_box(slide, 0.85, 3.85, 11.9, 1.75, "CVE_TECHNICAL_ANALYSIS", font_size=10, color=DARK_TEXT)
    _style_textbox(cta, DARK_TEXT, font_size=10)

    _add_card(slide, 0.65, 5.85, 12.3, 1.35, DARK_PANEL, DARK_BORDER)
    pt = add_placeholder_box(slide, 0.85, 5.95, 11.9, 0.4, "PATCH_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(pt, DARK_ACCENT, font_size=14, bold=True)
    ps = add_placeholder_box(slide, 0.85, 6.35, 11.9, 0.75, "PATCH_STATUS", font_size=10, color=DARK_TEXT)
    _style_textbox(ps, DARK_TEXT, font_size=10)

    # Slide 7: Attack Surface
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_base(slide)
    st = add_placeholder_box(slide, 0.65, 0.22, 12.3, 0.68, "SLIDE_TITLE", font_size=28, bold=True, color=DARK_ACCENT)
    _style_textbox(st, DARK_TEXT, font_size=28, bold=True)

    _add_card(slide, 0.65, 1.05, 6.05, 2.1, DARK_PANEL, DARK_BORDER)
    est = add_placeholder_box(slide, 0.85, 1.15, 5.65, 0.4, "EXPOSED_SERVICES_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(est, DARK_ACCENT, font_size=14, bold=True)
    estbl = add_placeholder_box(slide, 0.85, 1.55, 5.65, 1.5, "EXPOSED_SERVICES_TABLE", font_size=9, color=DARK_TEXT)
    _style_textbox(estbl, DARK_TEXT, font_size=9)

    _add_card(slide, 6.9, 1.05, 6.05, 2.1, DARK_PANEL, DARK_BORDER)
    eat = add_placeholder_box(slide, 7.1, 1.15, 5.65, 0.4, "EXPOSURE_ANALYSIS_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(eat, DARK_ACCENT, font_size=14, bold=True)
    ee = add_placeholder_box(slide, 7.1, 1.55, 5.65, 1.5, "EXTERNAL_EXPOSURE", font_size=10, color=DARK_TEXT)
    _style_textbox(ee, DARK_TEXT, font_size=10)

    _add_card(slide, 0.65, 3.3, 6.05, 3.9, DARK_PANEL, DARK_BORDER)
    stt = add_placeholder_box(slide, 0.85, 3.4, 5.65, 0.4, "SURFACE_TREND_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(stt, DARK_ACCENT, font_size=14, bold=True)
    ast = add_placeholder_box(slide, 0.85, 3.8, 5.65, 3.3, "ATTACK_SURFACE_TREND", font_size=11, color=DARK_TEXT)
    _style_textbox(ast, DARK_TEXT, font_size=11)

    _add_card(slide, 6.9, 3.3, 6.05, 3.9, DARK_PANEL, DARK_BORDER)
    rt = add_placeholder_box(slide, 7.1, 3.4, 5.65, 0.4, "REMEDIATION_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(rt, DARK_ACCENT, font_size=14, bold=True)
    er = add_placeholder_box(slide, 7.1, 3.8, 5.65, 3.3, "EXPOSURE_REMEDIATION", font_size=10, color=DARK_TEXT)
    _style_textbox(er, DARK_TEXT, font_size=10)

    # Slide 8: Cloud Security Details
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_base(slide)
    st = add_placeholder_box(slide, 0.65, 0.22, 12.3, 0.68, "SLIDE_TITLE", font_size=28, bold=True, color=DARK_ACCENT)
    _style_textbox(st, DARK_TEXT, font_size=28, bold=True)

    _add_card(slide, 0.65, 1.05, 12.3, 2.7, DARK_PANEL, DARK_BORDER)
    cmt = add_placeholder_box(slide, 0.85, 1.15, 11.9, 0.4, "CLOUD_MISCONFIG_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(cmt, DARK_ACCENT, font_size=14, bold=True)
    cmd = add_placeholder_box(slide, 0.85, 1.55, 11.9, 2.05, "CLOUD_MISCONFIG_DETAILS", font_size=10, color=DARK_TEXT)
    _style_textbox(cmd, DARK_TEXT, font_size=10)

    _add_card(slide, 0.65, 3.95, 6.05, 3.25, DARK_PANEL, DARK_BORDER)
    it = add_placeholder_box(slide, 0.85, 4.05, 5.65, 0.4, "IAM_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(it, DARK_ACCENT, font_size=14, bold=True)
    ira = add_placeholder_box(slide, 0.85, 4.45, 5.65, 2.65, "IAM_RISK_ANALYSIS", font_size=10, color=DARK_TEXT)
    _style_textbox(ira, DARK_TEXT, font_size=10)

    _add_card(slide, 6.9, 3.95, 6.05, 3.25, DARK_PANEL, DARK_BORDER)
    ct = add_placeholder_box(slide, 7.1, 4.05, 5.65, 0.4, "COMPLIANCE_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(ct, DARK_ACCENT, font_size=14, bold=True)
    cc = add_placeholder_box(slide, 7.1, 4.45, 5.65, 2.65, "CLOUD_COMPLIANCE", font_size=10, color=DARK_TEXT)
    _style_textbox(cc, DARK_TEXT, font_size=10)

    # Slide 9: Technical Recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_base(slide)
    st = add_placeholder_box(slide, 0.65, 0.22, 12.3, 0.68, "SLIDE_TITLE", font_size=28, bold=True, color=DARK_ACCENT)
    _style_textbox(st, DARK_TEXT, font_size=28, bold=True)

    _add_card(slide, 0.65, 1.05, 12.3, 2.45, DARK_PANEL, DARK_BORDER)
    trt = add_placeholder_box(slide, 0.85, 1.15, 11.9, 0.4, "TECH_REC_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(trt, DARK_ACCENT, font_size=14, bold=True)
    tr = add_placeholder_box(slide, 0.85, 1.55, 11.9, 1.85, "TECHNICAL_RECOMMENDATIONS", font_size=10, color=DARK_TEXT)
    _style_textbox(tr, DARK_TEXT, font_size=10)

    _add_card(slide, 0.65, 3.65, 6.05, 3.55, DARK_PANEL, DARK_BORDER)
    dt = add_placeholder_box(slide, 0.85, 3.75, 5.65, 0.4, "DETECTION_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(dt, DARK_ACCENT, font_size=14, bold=True)
    di = add_placeholder_box(slide, 0.85, 4.15, 5.65, 2.95, "DETECTION_IMPROVEMENTS", font_size=10, color=DARK_TEXT)
    _style_textbox(di, DARK_TEXT, font_size=10)

    _add_card(slide, 6.9, 3.65, 6.05, 3.55, DARK_PANEL, DARK_BORDER)
    ht = add_placeholder_box(slide, 7.1, 3.75, 5.65, 0.4, "HARDENING_TITLE", font_size=14, bold=True, color=DARK_ACCENT)
    _style_textbox(ht, DARK_ACCENT, font_size=14, bold=True)
    hc = add_placeholder_box(slide, 7.1, 4.15, 5.65, 2.95, "HARDENING_CHECKLIST", font_size=10, color=DARK_TEXT)
    _style_textbox(hc, DARK_TEXT, font_size=10)

    # Slide 10: Appendix
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_base(slide)
    st = add_placeholder_box(slide, 0.65, 0.22, 12.3, 0.68, "SLIDE_TITLE", font_size=28, bold=True, color=DARK_ACCENT)
    _style_textbox(st, DARK_TEXT, font_size=28, bold=True)

    _add_card(slide, 0.65, 1.05, 6.05, 1.2, DARK_PANEL, DARK_BORDER)
    et = add_placeholder_box(slide, 0.85, 1.15, 5.65, 0.35, "EVIDENCE_TITLE", font_size=12, bold=True, color=DARK_ACCENT)
    _style_textbox(et, DARK_ACCENT, font_size=12, bold=True)
    ei = add_placeholder_box(slide, 0.85, 1.45, 5.65, 0.75, "EVIDENCE_INDEX", font_size=9, color=DARK_TEXT)
    _style_textbox(ei, DARK_TEXT, font_size=9)

    _add_card(slide, 6.9, 1.05, 6.05, 1.2, DARK_PANEL, DARK_BORDER)
    ft = add_placeholder_box(slide, 7.1, 1.15, 5.65, 0.35, "FRESHNESS_TITLE", font_size=12, bold=True, color=DARK_ACCENT)
    _style_textbox(ft, DARK_ACCENT, font_size=12, bold=True)
    df = add_placeholder_box(slide, 7.1, 1.45, 5.65, 0.75, "DATA_FRESHNESS", font_size=10, color=DARK_TEXT)
    _style_textbox(df, DARK_TEXT, font_size=10)

    _add_card(slide, 0.65, 2.35, 12.3, 0.85, DARK_PANEL, DARK_BORDER)
    dst = add_placeholder_box(slide, 0.85, 2.45, 3.4, 0.35, "DATA_SCOPE_TITLE", font_size=12, bold=True, color=DARK_ACCENT)
    _style_textbox(dst, DARK_ACCENT, font_size=12, bold=True)
    dsn = add_placeholder_box(slide, 4.25, 2.45, 8.7, 0.55, "DATA_SCOPE_NOTES", font_size=10, color=DARK_TEXT)
    _style_textbox(dsn, DARK_TEXT, font_size=10)

    _add_card(slide, 0.65, 3.3, 6.05, 1.35, DARK_PANEL, DARK_BORDER)
    vnt = add_placeholder_box(slide, 0.85, 3.4, 5.65, 0.35, "VULN_NOTES_TITLE", font_size=12, bold=True, color=DARK_ACCENT)
    _style_textbox(vnt, DARK_ACCENT, font_size=12, bold=True)
    vn = add_placeholder_box(slide, 0.85, 3.7, 5.65, 0.85, "VULN_NOTES", font_size=10, color=DARK_TEXT)
    _style_textbox(vn, DARK_TEXT, font_size=10)

    _add_card(slide, 6.9, 3.3, 6.05, 1.35, DARK_PANEL, DARK_BORDER)
    ent = add_placeholder_box(slide, 7.1, 3.4, 5.65, 0.35, "EVIDENCE_NOTES_TITLE", font_size=12, bold=True, color=DARK_ACCENT)
    _style_textbox(ent, DARK_ACCENT, font_size=12, bold=True)
    en = add_placeholder_box(slide, 7.1, 3.7, 5.65, 0.85, "EVIDENCE_NOTES", font_size=10, color=DARK_TEXT)
    _style_textbox(en, DARK_TEXT, font_size=10)

    _add_card(slide, 0.65, 4.75, 6.05, 1.15, DARK_PANEL, DARK_BORDER)
    at = add_placeholder_box(slide, 0.85, 4.85, 5.65, 0.35, "ASSETS_TITLE", font_size=12, bold=True, color=DARK_ACCENT)
    _style_textbox(at, DARK_ACCENT, font_size=12, bold=True)
    ag = add_placeholder_box(slide, 0.85, 5.15, 5.65, 0.65, "ASSET_GROUPS_LIST", font_size=10, color=DARK_TEXT)
    _style_textbox(ag, DARK_TEXT, font_size=10)

    _add_card(slide, 6.9, 4.75, 6.05, 1.15, DARK_PANEL, DARK_BORDER)
    svt = add_placeholder_box(slide, 7.1, 4.85, 5.65, 0.35, "SERVICES_TITLE", font_size=12, bold=True, color=DARK_ACCENT)
    _style_textbox(svt, DARK_ACCENT, font_size=12, bold=True)
    ksl = add_placeholder_box(slide, 7.1, 5.15, 5.65, 0.65, "KEY_SERVICES_LIST", font_size=10, color=DARK_TEXT)
    _style_textbox(ksl, DARK_TEXT, font_size=10)

    _add_card(slide, 0.65, 6.0, 12.3, 1.45, DARK_PANEL, DARK_BORDER)
    termt = add_placeholder_box(slide, 0.85, 6.1, 11.9, 0.35, "TERMINOLOGY_TITLE", font_size=12, bold=True, color=DARK_ACCENT)
    _style_textbox(termt, DARK_ACCENT, font_size=12, bold=True)
    term = add_placeholder_box(slide, 0.85, 6.4, 11.9, 1.0, "TERMINOLOGY_TECHNICAL", font_size=9, color=DARK_TEXT)
    _style_textbox(term, DARK_TEXT, font_size=9)

    return prs


def main():
    """Generate both V2 templates."""
    output_dir = Path(__file__).parent / "data" / "templates"
    output_dir.mkdir(parents=True, exist_ok=True)

    # Generate executive template
    print("Generating mss_executive_v2.pptx...")
    exec_prs = create_executive_template()
    exec_path = output_dir / "mss_executive_v2.pptx"
    exec_prs.save(exec_path)
    print(f"  Saved to: {exec_path}")

    # Generate technical template
    print("Generating mss_technical_v2.pptx...")
    tech_prs = create_technical_template()
    tech_path = output_dir / "mss_technical_v2.pptx"
    tech_prs.save(tech_path)
    print(f"  Saved to: {tech_path}")

    print("\nDone! V2 templates generated successfully.")


if __name__ == "__main__":
    main()
