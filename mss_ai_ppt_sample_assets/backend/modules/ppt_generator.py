from __future__ import annotations

import logging
from pathlib import Path
from typing import Dict, Any, List, Optional

from mss_ai_ppt_sample_assets.backend.models.slidespec import SlideSpecV2
from mss_ai_ppt_sample_assets.backend.modules.template_loader import TemplateRepository

logger = logging.getLogger(__name__)


# Professional color palettes for charts and tables
class ChartColors:
    """Professional color schemes for charts and tables."""

    # Primary palette (blue-based, professional)
    PRIMARY = [
        (30, 64, 175),    # Deep Blue
        (59, 130, 246),   # Bright Blue
        (147, 197, 253),  # Light Blue
        (96, 165, 250),   # Medium Blue
    ]

    # Severity palette (for security data)
    SEVERITY = {
        'critical': (220, 38, 38),   # Red
        'high': (234, 88, 12),       # Orange
        'medium': (234, 179, 8),     # Yellow
        'low': (34, 197, 94),        # Green
        'info': (148, 163, 184),     # Gray
    }

    # Multi-series palette (vibrant, distinguishable)
    MULTI_SERIES = [
        (30, 64, 175),    # Deep Blue
        (34, 197, 94),    # Green
        (234, 179, 8),    # Yellow/Amber
        (168, 85, 247),   # Purple
        (236, 72, 153),   # Pink
        (20, 184, 166),   # Teal
    ]

    # Table styles
    TABLE_HEADER_BG = (30, 64, 175)       # Deep Blue
    TABLE_HEADER_TEXT = (255, 255, 255)   # White
    TABLE_ROW_ALT = (241, 245, 249)       # Light Gray (alternating)
    TABLE_ROW_NORMAL = (255, 255, 255)    # White
    TABLE_BORDER = (203, 213, 225)        # Slate Gray

    # Theme-aware text colors
    TEXT_LIGHT_THEME = (30, 41, 59)       # Slate-800 for light backgrounds
    TEXT_DARK_THEME = (255, 255, 255)     # White for dark backgrounds

    # Dark theme table styles
    TABLE_HEADER_BG_DARK = (34, 197, 94)  # Green accent for dark theme
    TABLE_ROW_ALT_DARK = (30, 41, 59)     # Slate-800 (alternating)
    TABLE_ROW_NORMAL_DARK = (15, 23, 42)  # Slate-900


class PPTGeneratorV2:
    """Fill PPTX template placeholders with V2 slidespec content.

    V2 simplification: slidespec.placeholders directly maps token -> value,
    no need to traverse render_map paths.
    """

    def __init__(self, template_repo: TemplateRepository):
        self.template_repo = template_repo
        self._is_dark_theme = False  # Track current template theme
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
            from pptx.chart.data import CategoryChartData
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.dml.color import RGBColor

            self._Presentation = Presentation
            self._Inches = Inches
            self._Pt = Pt
            self._Emu = Emu
            self._XL_CHART_TYPE = XL_CHART_TYPE
            self._XL_LEGEND_POSITION = XL_LEGEND_POSITION
            self._XL_LABEL_POSITION = XL_LABEL_POSITION
            self._CategoryChartData = CategoryChartData
            self._PP_ALIGN = PP_ALIGN
            self._MSO_ANCHOR = MSO_ANCHOR
            self._RGBColor = RGBColor
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx is required for PPT rendering. Please install via requirements.txt."
            ) from exc

    def _replace_tokens_in_shape(self, shape, mapping: Dict[str, str]) -> None:
        """Replace {{TOKEN}} placeholders in shape text."""
        if not shape.has_text_frame:
            return
        
        shape_modified = False
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                for token, value in mapping.items():
                    placeholder = f"{{{{{token}}}}}"
                    if placeholder in text:
                        run.text = text.replace(placeholder, str(value) if value else "")
                        text = run.text
                        shape_modified = True
            
            # Apply formatting if modified and text is substantial (multiline or long)
            # This helps avoid "crowded" look for generated content
            if shape_modified:
                # Re-check text content after replacement
                p_text = paragraph.text
                if '\n' in p_text or len(p_text) > 50:
                    paragraph.line_spacing = 1.25
                    paragraph.space_after = self._Pt(6)
                    # Ensure wrapping is on for long text
                    if shape.text_frame.word_wrap is None:
                        shape.text_frame.word_wrap = True

    def _render_native_table(
        self,
        slide,
        table_data: Dict[str, Any],
        position: Optional[Dict[str, float]] = None
    ) -> None:
        """Render a professional-styled native PowerPoint table.

        Args:
            slide: pptx slide object
            table_data: Dict with 'headers' (list) and 'rows' (list of lists)
            position: Dict with 'left', 'top', 'width', 'height' in inches
        """
        if not table_data or 'headers' not in table_data or 'rows' not in table_data:
            logger.warning("Invalid table data format")
            return

        headers = table_data['headers']
        rows = table_data['rows']

        if not headers or not rows:
            logger.warning("Empty table data")
            return

        # Default position if not specified
        if position is None:
            position = {'left': 0.8, 'top': 2.5, 'width': 8.5, 'height': 3.0}

        # Calculate dimensions
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)

        # Create table
        left = self._Inches(position.get('left', 0.8))
        top = self._Inches(position.get('top', 2.5))
        width = self._Inches(position.get('width', 8.5))
        height = self._Inches(position.get('height', 3.0))

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        # Calculate column widths based on header lengths and content
        col_widths = table_data.get('col_widths', None)
        if col_widths:
            total_width = sum(col_widths)
            for col_idx, cw in enumerate(col_widths):
                table.columns[col_idx].width = self._Inches(cw * position.get('width', 8.5) / total_width)

        # Style headers (professional deep blue with white text)
        header_bg = ChartColors.TABLE_HEADER_BG
        header_text = ChartColors.TABLE_HEADER_TEXT

        for col_idx, header in enumerate(headers):
            cell = table.rows[0].cells[col_idx]
            cell.text = str(header)

            # Header background
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._RGBColor(*header_bg)

            # Header text styling
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = self._Pt(11)
            paragraph.font.color.rgb = self._RGBColor(*header_text)
            paragraph.font.name = "微软雅黑"
            paragraph.alignment = self._PP_ALIGN.CENTER

            # Vertical alignment
            cell.vertical_anchor = self._MSO_ANCHOR.MIDDLE

        # Style data rows with alternating colors
        for row_idx, row_data in enumerate(rows):
            # Alternating row colors for better readability
            if row_idx % 2 == 0:
                row_bg = ChartColors.TABLE_ROW_NORMAL
            else:
                row_bg = ChartColors.TABLE_ROW_ALT

            for col_idx, cell_value in enumerate(row_data):
                if col_idx < num_cols:
                    cell = table.rows[row_idx + 1].cells[col_idx]

                    # Format cell value
                    display_value = str(cell_value) if cell_value is not None else ""
                    cell.text = display_value

                    # Row background
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._RGBColor(*row_bg)

                    # Data cell text styling
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = self._Pt(10)
                    paragraph.font.name = "微软雅黑"
                    paragraph.font.color.rgb = self._RGBColor(30, 41, 59)  # Slate-800

                    # Center numeric columns, left-align text
                    if isinstance(cell_value, (int, float)) or (isinstance(cell_value, str) and cell_value.replace('.', '').replace('%', '').isdigit()):
                        paragraph.alignment = self._PP_ALIGN.CENTER
                    else:
                        paragraph.alignment = self._PP_ALIGN.LEFT

                    # Vertical alignment
                    cell.vertical_anchor = self._MSO_ANCHOR.MIDDLE

        logger.info(f"Rendered professional table: {num_rows} rows x {num_cols} cols")

    def _render_bar_chart(
        self,
        slide,
        chart_data: Dict[str, Any],
        position: Optional[Dict[str, float]] = None
    ) -> None:
        """Render a professional bar/column chart with enhanced styling.

        Args:
            slide: pptx slide object
            chart_data: Dict with 'categories' (list), 'series' (list of dicts with 'name' and 'values')
            position: Dict with 'left', 'top', 'width', 'height' in inches
        """
        if not chart_data or 'categories' not in chart_data or 'series' not in chart_data:
            logger.warning("Invalid bar chart data format")
            return

        categories = chart_data['categories']
        series_list = chart_data['series']

        if not categories or not series_list:
            logger.warning("Empty bar chart data")
            return

        # Default position if not specified
        if position is None:
            position = {'left': 0.8, 'top': 2.0, 'width': 8.5, 'height': 4.5}

        # Create chart data
        chart_data_obj = self._CategoryChartData()
        chart_data_obj.categories = categories

        for series in series_list:
            series_name = series.get('name', 'Series')
            series_values = series.get('values', [])
            chart_data_obj.add_series(series_name, series_values)

        # Add chart to slide
        left = self._Inches(position.get('left', 0.8))
        top = self._Inches(position.get('top', 2.0))
        width = self._Inches(position.get('width', 8.5))
        height = self._Inches(position.get('height', 4.5))

        chart_shape = slide.shapes.add_chart(
            self._XL_CHART_TYPE.COLUMN_CLUSTERED,
            left, top, width, height,
            chart_data_obj
        )
        chart = chart_shape.chart

        # Use theme-aware text color
        text_color = ChartColors.TEXT_DARK_THEME if self._is_dark_theme else ChartColors.TEXT_LIGHT_THEME
        axis_text_color = ChartColors.TEXT_DARK_THEME if self._is_dark_theme else (71, 85, 105)  # Slate-500 for light theme

        # Enhanced styling
        # Set chart title if provided
        if 'title' in chart_data and chart_data['title']:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_data['title']
            chart.chart_title.text_frame.paragraphs[0].font.size = self._Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self._RGBColor(*text_color)

        # Style the series with professional colors
        for idx, series in enumerate(chart.series):
            color = ChartColors.MULTI_SERIES[idx % len(ChartColors.MULTI_SERIES)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = self._RGBColor(*color)

            # Add data labels
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.font.size = self._Pt(9)
            data_labels.font.color.rgb = self._RGBColor(*text_color)
            data_labels.number_format = '#,##0'

        # Configure legend
        chart.has_legend = len(series_list) > 1
        if chart.has_legend:
            chart.legend.position = self._XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = self._Pt(10)
            chart.legend.font.color.rgb = self._RGBColor(*text_color)

        # Style category axis
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = self._Pt(10)
        category_axis.tick_labels.font.color.rgb = self._RGBColor(*axis_text_color)

        # Style value axis
        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = self._Pt(9)
        value_axis.tick_labels.font.color.rgb = self._RGBColor(*axis_text_color)
        value_axis.has_major_gridlines = True

        logger.info(f"Rendered professional bar chart with {len(categories)} categories and {len(series_list)} series")

    def _render_pie_chart(
        self,
        slide,
        chart_data: Dict[str, Any],
        position: Optional[Dict[str, float]] = None
    ) -> None:
        """Render a professional pie chart with enhanced styling.

        Args:
            slide: pptx slide object
            chart_data: Dict with 'categories' (list) and 'values' (list)
            position: Dict with 'left', 'top', 'width', 'height' in inches
        """
        if not chart_data or 'categories' not in chart_data or 'values' not in chart_data:
            logger.warning("Invalid pie chart data format")
            return

        categories = chart_data['categories']
        values = chart_data['values']

        if not categories or not values or len(categories) != len(values):
            logger.warning("Invalid or mismatched pie chart data")
            return

        # Default position if not specified
        if position is None:
            position = {'left': 2.0, 'top': 2.0, 'width': 6.0, 'height': 4.5}

        # Create chart data
        chart_data_obj = self._CategoryChartData()
        chart_data_obj.categories = categories
        chart_data_obj.add_series('', values)

        # Add chart to slide
        left = self._Inches(position.get('left', 2.0))
        top = self._Inches(position.get('top', 2.0))
        width = self._Inches(position.get('width', 6.0))
        height = self._Inches(position.get('height', 4.5))

        chart_shape = slide.shapes.add_chart(
            self._XL_CHART_TYPE.PIE,
            left, top, width, height,
            chart_data_obj
        )
        chart = chart_shape.chart

        # Use theme-aware text color
        text_color = ChartColors.TEXT_DARK_THEME if self._is_dark_theme else ChartColors.TEXT_LIGHT_THEME

        # Enhanced styling
        # Set chart title if provided
        if 'title' in chart_data and chart_data['title']:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_data['title']
            chart.chart_title.text_frame.paragraphs[0].font.size = self._Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
            chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self._RGBColor(*text_color)

        # Get color scheme - use severity colors if categories match severity levels
        severity_keywords = ['严重', '高危', '中危', '低危', '信息', 'critical', 'high', 'medium', 'low', 'info']
        use_severity_colors = any(kw in str(cat).lower() for cat in categories for kw in severity_keywords)

        # Style pie slices with professional colors
        plot = chart.plots[0]
        for idx, point in enumerate(plot.series[0].points):
            if use_severity_colors:
                # Map category to severity color
                cat_lower = str(categories[idx]).lower()
                if '严重' in cat_lower or 'critical' in cat_lower:
                    color = ChartColors.SEVERITY['critical']
                elif '高危' in cat_lower or 'high' in cat_lower:
                    color = ChartColors.SEVERITY['high']
                elif '中危' in cat_lower or 'medium' in cat_lower:
                    color = ChartColors.SEVERITY['medium']
                elif '低危' in cat_lower or 'low' in cat_lower:
                    color = ChartColors.SEVERITY['low']
                else:
                    color = ChartColors.SEVERITY['info']
            else:
                color = ChartColors.MULTI_SERIES[idx % len(ChartColors.MULTI_SERIES)]

            point.format.fill.solid()
            point.format.fill.fore_color.rgb = self._RGBColor(*color)

        # Add data labels with percentages
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_category_name = True
        data_labels.show_percentage = True
        data_labels.show_value = False
        data_labels.font.size = self._Pt(10)
        data_labels.font.color.rgb = self._RGBColor(*text_color)
        data_labels.number_format = '0.0%'

        # Configure legend
        chart.has_legend = True
        chart.legend.position = self._XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = self._Pt(10)
        chart.legend.font.color.rgb = self._RGBColor(*text_color)

        logger.info(f"Rendered professional pie chart with {len(categories)} categories")

    def _process_chart_placeholder(
        self,
        slide,
        token: str,
        value: Any,
        chart_type: str
    ) -> bool:
        """Process a chart placeholder value.

        Args:
            slide: pptx slide object
            token: placeholder token name
            value: chart data (dict or special format)
            chart_type: 'bar_chart' or 'pie_chart'

        Returns:
            True if chart was rendered, False otherwise
        """
        if not isinstance(value, dict):
            logger.warning(f"Chart placeholder {token} has invalid data type: {type(value)}")
            return False

        try:
            if chart_type == 'bar_chart':
                self._render_bar_chart(slide, value, value.get('position'))
            elif chart_type == 'pie_chart':
                self._render_pie_chart(slide, value, value.get('position'))
            else:
                logger.warning(f"Unknown chart type: {chart_type}")
                return False
            return True
        except Exception as e:
            logger.error(f"Failed to render {chart_type} for {token}: {e}")
            return False

    def _process_table_placeholder(
        self,
        slide,
        token: str,
        value: Any
    ) -> bool:
        """Process a native table placeholder value.

        Args:
            slide: pptx slide object
            token: placeholder token name
            value: table data (dict with headers and rows)

        Returns:
            True if table was rendered, False otherwise
        """
        if not isinstance(value, dict):
            logger.warning(f"Table placeholder {token} has invalid data type: {type(value)}")
            return False

        try:
            self._render_native_table(slide, value, value.get('position'))
            return True
        except Exception as e:
            logger.error(f"Failed to render table for {token}: {e}")
            return False

    def render(self, slidespec: SlideSpecV2, output_path: Path) -> Path:
        """Render V2 slidespec to PPTX.

        Args:
            slidespec: V2 slidespec with placeholders filled
            output_path: Where to save the output PPTX

        Returns:
            Path to the saved PPTX file
        """
        template_path = self.template_repo.get_pptx_path(slidespec.template_id)
        prs = self._Presentation(template_path)

        # Load template descriptor to get placeholder types
        template_desc = self.template_repo.get_descriptor_v2(slidespec.template_id)

        # Detect theme from template descriptor style
        self._is_dark_theme = template_desc.style.get('theme', 'light') == 'dark'

        # Build mapping: slide_key -> placeholder_definitions
        placeholder_types = {}
        for slide_def in template_desc.slides:
            placeholder_types[slide_def.slide_key] = {
                ph.token: ph.type for ph in slide_def.placeholders
            }

        # Build mapping: slide_no -> placeholders dict
        slides_by_no = {s.slide_no: s for s in slidespec.slides}

        for slide_no, slide_content in slides_by_no.items():
            if slide_no > len(prs.slides):
                continue

            # pptx slides are 0-indexed
            pptx_slide = prs.slides[slide_no - 1]

            # Get placeholder types for this slide
            slide_types = placeholder_types.get(slide_content.slide_key, {})

            # Separate placeholders by type
            text_placeholders = {}
            chart_bar_placeholders = []
            chart_pie_placeholders = []
            table_placeholders = []

            for token, value in slide_content.placeholders.items():
                ph_type = slide_types.get(token, 'text')

                if ph_type in ('bar_chart',):
                    chart_bar_placeholders.append((token, value))
                elif ph_type in ('pie_chart',):
                    chart_pie_placeholders.append((token, value))
                elif ph_type in ('native_table',):
                    table_placeholders.append((token, value))
                else:
                    # Regular text placeholder
                    if value is None:
                        value = ""
                    elif isinstance(value, list):
                        # Join list items with newlines
                        value = "\n".join(str(v) for v in value)
                    text_placeholders[token] = str(value)

            # Replace text tokens in all shapes
            for shape in pptx_slide.shapes:
                self._replace_tokens_in_shape(shape, text_placeholders)

            # Render charts
            for token, value in chart_bar_placeholders:
                self._process_chart_placeholder(pptx_slide, token, value, 'bar_chart')

            for token, value in chart_pie_placeholders:
                self._process_chart_placeholder(pptx_slide, token, value, 'pie_chart')

            # Render tables
            for token, value in table_placeholders:
                self._process_table_placeholder(pptx_slide, token, value)

        # Save output
        output_path.parent.mkdir(parents=True, exist_ok=True)
        if output_path.exists():
            try:
                output_path.unlink()
            except OSError:
                output_path = output_path.with_name(output_path.stem + "_new" + output_path.suffix)

        prs.save(output_path)
        return output_path